#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
TX-105: ブランディング戦略資料（PowerPoint）- 検収物独自性欠如
汎用テンプレートそのままの成果物を生成
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN
import os

def create_generic_branding_ppt():
    """汎用的なブランディング戦略テンプレート（独自性なし）"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # スライド1: タイトル
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RgbColor(31, 78, 121)  # 濃紺
    
    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "ブランディング戦略提案書"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # サブタイトル
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "〜 貴社のブランド価値向上に向けて 〜"
    p.font.size = Pt(24)
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 会社名
    company_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.6))
    tf = company_box.text_frame
    p = tf.paragraphs[0]
    p.text = "株式会社ストラテジーパートナーズ"
    p.font.size = Pt(18)
    p.font.color.rgb = RgbColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER
    
    # スライド2: 目次（汎用的）
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "目次"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RgbColor(31, 78, 121)
    
    contents = [
        "1. 現状分析",
        "2. 課題の整理",
        "3. ブランド戦略の方向性",
        "4. 具体的施策",
        "5. 実行計画",
        "6. 期待効果"
    ]
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(5))
    tf = content_box.text_frame
    for i, item in enumerate(contents):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(24)
        p.font.color.rgb = RgbColor(50, 50, 50)
        p.space_after = Pt(20)
    
    # スライド3: 現状分析（汎用テンプレート）
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "1. 現状分析"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(31, 78, 121)
    
    # 汎用的な分析フレームワーク（SWOT風）
    boxes = [
        (Inches(0.5), Inches(1.5), "強み (Strengths)", "・〇〇〇〇〇〇\n・〇〇〇〇〇〇\n・〇〇〇〇〇〇"),
        (Inches(6.5), Inches(1.5), "弱み (Weaknesses)", "・〇〇〇〇〇〇\n・〇〇〇〇〇〇\n・〇〇〇〇〇〇"),
        (Inches(0.5), Inches(4.2), "機会 (Opportunities)", "・〇〇〇〇〇〇\n・〇〇〇〇〇〇\n・〇〇〇〇〇〇"),
        (Inches(6.5), Inches(4.2), "脅威 (Threats)", "・〇〇〇〇〇〇\n・〇〇〇〇〇〇\n・〇〇〇〇〇〇"),
    ]
    
    for x, y, title, content in boxes:
        # ボックス
        shape = slide.shapes.add_shape(1, x, y, Inches(5.8), Inches(2.5))  # Rectangle
        shape.fill.solid()
        shape.fill.fore_color.rgb = RgbColor(240, 240, 240)
        shape.line.color.rgb = RgbColor(31, 78, 121)
        
        # タイトル
        title_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.1), Inches(5.4), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RgbColor(31, 78, 121)
        
        # 内容（プレースホルダー）
        content_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.6), Inches(5.4), Inches(1.8))
        tf = content_box.text_frame
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(14)
        p.font.color.rgb = RgbColor(100, 100, 100)
    
    # スライド4: 課題の整理（汎用）
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "2. 課題の整理"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(31, 78, 121)
    
    issues = [
        "課題1: ブランド認知度の向上が必要",
        "課題2: 競合との差別化ポイントの明確化",
        "課題3: ターゲット顧客への効果的なリーチ",
        "課題4: 一貫したブランドメッセージの発信"
    ]
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11), Inches(5))
    tf = content_box.text_frame
    for i, issue in enumerate(issues):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = issue
        p.font.size = Pt(22)
        p.font.color.rgb = RgbColor(50, 50, 50)
        p.space_after = Pt(30)
    
    # スライド5: ブランド戦略の方向性（汎用）
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "3. ブランド戦略の方向性"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(31, 78, 121)
    
    # 汎用的な戦略図
    strategies = [
        (Inches(1), "ブランドビジョン", "〇〇〇〇〇〇〇〇〇〇"),
        (Inches(3.5), "ブランドミッション", "〇〇〇〇〇〇〇〇〇〇"),
        (Inches(6), "ブランドバリュー", "〇〇〇〇〇〇〇〇〇〇"),
    ]
    
    for y, title, content in strategies:
        shape = slide.shapes.add_shape(1, Inches(1), y, Inches(11), Inches(1.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RgbColor(31, 78, 121)
        
        title_box = slide.shapes.add_textbox(Inches(1.3), y + Inches(0.2), Inches(10), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RgbColor(255, 255, 255)
        
        content_box = slide.shapes.add_textbox(Inches(1.3), y + Inches(0.8), Inches(10), Inches(0.8))
        tf = content_box.text_frame
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(16)
        p.font.color.rgb = RgbColor(200, 200, 200)
    
    # スライド6: 具体的施策（汎用）
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "4. 具体的施策"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(31, 78, 121)
    
    measures = [
        ("施策1", "〇〇〇〇〇〇〇〇"),
        ("施策2", "〇〇〇〇〇〇〇〇"),
        ("施策3", "〇〇〇〇〇〇〇〇"),
    ]
    
    for i, (title, content) in enumerate(measures):
        x = Inches(0.5 + i * 4.2)
        shape = slide.shapes.add_shape(1, x, Inches(1.5), Inches(3.8), Inches(4.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RgbColor(240, 240, 240)
        shape.line.color.rgb = RgbColor(31, 78, 121)
        
        title_box = slide.shapes.add_textbox(x + Inches(0.2), Inches(1.7), Inches(3.4), Inches(0.6))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RgbColor(31, 78, 121)
        p.alignment = PP_ALIGN.CENTER
        
        content_box = slide.shapes.add_textbox(x + Inches(0.2), Inches(2.5), Inches(3.4), Inches(3))
        tf = content_box.text_frame
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(14)
        p.font.color.rgb = RgbColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER
    
    # スライド7: 実行計画（汎用ガントチャート風）
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "5. 実行計画"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(31, 78, 121)
    
    # 簡易ガントチャート
    phases = ["Phase 1", "Phase 2", "Phase 3"]
    months = ["1月", "2月", "3月", "4月", "5月", "6月"]
    
    # ヘッダー
    for i, month in enumerate(months):
        box = slide.shapes.add_textbox(Inches(3 + i * 1.5), Inches(1.5), Inches(1.5), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = month
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.CENTER
    
    # フェーズ
    for i, phase in enumerate(phases):
        y = Inches(2.2 + i * 1.2)
        box = slide.shapes.add_textbox(Inches(0.5), y, Inches(2.3), Inches(0.8))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = phase
        p.font.size = Pt(14)
        p.font.bold = True
        
        # バー
        bar_start = 3 + i * 1.5
        bar_width = 3 - i * 0.5
        shape = slide.shapes.add_shape(1, Inches(bar_start), y + Inches(0.1), Inches(bar_width), Inches(0.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RgbColor(31, 78, 121)
    
    # スライド8: 期待効果（汎用）
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "6. 期待効果"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(31, 78, 121)
    
    effects = [
        "ブランド認知度: 〇〇%向上",
        "顧客ロイヤルティ: 〇〇%改善",
        "売上貢献: 〇〇%増加見込み"
    ]
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
    tf = content_box.text_frame
    for i, effect in enumerate(effects):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = effect
        p.font.size = Pt(28)
        p.font.color.rgb = RgbColor(50, 50, 50)
        p.space_after = Pt(40)
    
    # 保存
    output_path = '/home/ubuntu/ad_audit_test_data/検収資料/PowerPoint/TX-105_ブランディング戦略資料.pptx'
    prs.save(output_path)
    print(f"生成完了: {output_path}")


if __name__ == "__main__":
    create_generic_branding_ppt()
